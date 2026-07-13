from pptx import Presentation
from langchain_core.documents import Document

def extract_ppt_documents(ppt_path: str) -> str:
    """
    Extract text from PowerPoint.
    """

    prs = Presentation(ppt_path)

    documents = []

    for slide_number, slide in enumerate(prs.slides):
        slide_text = ""

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"

        documents.append(
            Document(
                page_content=slide_text,
                metadata={
                    "slide": slide_number+1,
                    "source": ppt_path
                }
            )
        )
    
    return documents